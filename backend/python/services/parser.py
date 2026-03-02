"""
services/parser.py
Handles text extraction from every supported source type:
  - PDF (PyMuPDF + embedded-image OCR)
  - Images (pytesseract OCR)
  - DOCX / PPTX
  - Plain text / TXT
  - YouTube links (transcript API)
  - Websites / documentation URLs (BeautifulSoup)
"""

import fitz                             # PyMuPDF
import pytesseract
from PIL import Image
import io, os, re, tempfile
from pathlib import Path
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
from youtube_transcript_api import YouTubeTranscriptApi
from models.schemas import SourceType


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _clean(text: str) -> str:
    """Remove excessive whitespace / blank lines."""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


def _youtube_id(url: str) -> str | None:
    """Extract YouTube video ID from various URL formats."""
    patterns = [
        r'youtube\.com/watch\?v=([A-Za-z0-9_-]{11})',
        r'youtu\.be/([A-Za-z0-9_-]{11})',
        r'youtube\.com/embed/([A-Za-z0-9_-]{11})',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


# ──────────────────────────────────────────────
# Per-type parsers
# ──────────────────────────────────────────────

def parse_pdf(file_path: str) -> list[dict]:
    """
    Returns list of {page: int, text: str} dicts.
    Also runs OCR on embedded images when text is sparse.
    """
    pages = []
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        # If page has very little text, try OCR on the rendered page image
        if len(text.strip()) < 50:
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img)
        pages.append({"page": page_num, "text": _clean(text)})
    doc.close()
    return [p for p in pages if p["text"]]


def parse_image(file_path: str) -> list[dict]:
    """OCR on image files (JPEG, PNG, WEBP, TIFF, BMP)."""
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return [{"page": 1, "text": _clean(text)}] if text.strip() else []


def parse_docx(file_path: str) -> list[dict]:
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    return [{"page": None, "text": _clean(text)}] if text else []


def parse_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append({"page": i, "text": _clean("\n".join(texts))})
    return slides


def parse_txt(file_path: str) -> list[dict]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [{"page": None, "text": _clean(text)}] if text.strip() else []


def parse_youtube(url: str) -> list[dict]:
    """Fetch transcript from YouTube."""
    video_id = _youtube_id(url)
    if not video_id:
        raise ValueError(f"Could not extract YouTube video ID from URL: {url}")
    transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=["en", "en-US"])
    # Group transcript into ~60 second chunks
    pages = []
    current, current_time = [], 0.0
    chunk_duration = 60  # seconds per chunk
    for entry in transcript_list:
        if entry["start"] - current_time > chunk_duration and current:
            pages.append({"page": len(pages) + 1, "text": _clean(" ".join(current))})
            current = []
            current_time = entry["start"]
        current.append(entry["text"])
    if current:
        pages.append({"page": len(pages) + 1, "text": _clean(" ".join(current))})
    return pages


def parse_website(url: str) -> list[dict]:
    """Scrape text from a website / documentation page."""
    headers = {"User-Agent": "Mozilla/5.0 (Mentoroid RAG Bot)"}
    resp = requests.get(url, headers=headers, timeout=20)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "lxml")
    # Remove script/style/nav junk
    for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form"]):
        tag.decompose()
    # Try to get <main> or <article> first for cleaner content
    main = soup.find("main") or soup.find("article") or soup.body
    text = main.get_text(separator="\n") if main else soup.get_text(separator="\n")
    return [{"page": None, "text": _clean(text)}] if text.strip() else []


# ──────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────

EXTENSION_MAP = {
    ".pdf":  SourceType.pdf,
    ".jpg":  SourceType.image,
    ".jpeg": SourceType.image,
    ".png":  SourceType.image,
    ".webp": SourceType.image,
    ".bmp":  SourceType.image,
    ".tiff": SourceType.image,
    ".tif":  SourceType.image,
    ".docx": SourceType.docx,
    ".doc":  SourceType.docx,
    ".pptx": SourceType.pptx,
    ".ppt":  SourceType.pptx,
    ".txt":  SourceType.txt,
    ".md":   SourceType.txt,
}


def detect_source_type(filename: str, url: str = None) -> SourceType:
    if url:
        if _youtube_id(url):
            return SourceType.youtube
        return SourceType.website
    ext = Path(filename).suffix.lower()
    return EXTENSION_MAP.get(ext, SourceType.txt)


def parse_file(file_path: str, source_type: SourceType) -> list[dict]:
    """Main dispatcher — returns list of {page, text} dicts."""
    dispatch = {
        SourceType.pdf:   parse_pdf,
        SourceType.image: parse_image,
        SourceType.docx:  parse_docx,
        SourceType.pptx:  parse_pptx,
        SourceType.txt:   parse_txt,
    }
    fn = dispatch.get(source_type)
    if fn is None:
        raise ValueError(f"parse_file() called with URL type '{source_type}' — use parse_url() instead")
    return fn(file_path)


def parse_url(url: str, source_type: SourceType) -> list[dict]:
    """Dispatcher for URL-based sources."""
    if source_type == SourceType.youtube:
        return parse_youtube(url)
    return parse_website(url)
